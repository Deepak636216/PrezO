"""
Enhanced Template Analysis Tool with Gemini AI
Extracts structure from PPTX templates and uses AI to generate functions and descriptions
"""

import os
import json
import base64
from io import BytesIO
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()


def generate_template_id(template_file_path: str) -> str:
    """Generate unique template ID from file path"""
    basename = os.path.basename(template_file_path)
    name_without_ext = os.path.splitext(basename)[0]
    template_id = name_without_ext.lower().replace(' ', '_').replace('-', '_')
    return template_id


def export_slide_as_image(slide, slide_prs) -> bytes:
    """
    Export a single slide as PNG image for AI analysis

    Args:
        slide: Slide object
        slide_prs: Presentation containing the slide

    Returns:
        PNG image bytes
    """
    # Create temporary presentation with just this slide
    temp_prs = Presentation()
    temp_prs.slide_width = slide_prs.slide_width
    temp_prs.slide_height = slide_prs.slide_height

    # Clone the slide
    blank_layout = temp_prs.slide_layouts[0]
    new_slide = temp_prs.slides.add_slide(blank_layout)

    # Copy all shapes
    for shape in slide.shapes:
        try:
            import copy
            el = shape.element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        except:
            pass

    # Save to BytesIO
    img_bytes = BytesIO()

    # Note: python-pptx doesn't support direct image export
    # We'll extract shape information instead
    return None


def analyze_slide_with_gemini(slide, slide_index: int, api_key: str) -> Dict[str, Any]:
    """
    Use Gemini to analyze slide and generate description

    Args:
        slide: Slide object
        slide_index: Index of the slide
        api_key: Google AI Studio API key

    Returns:
        Dictionary with AI-generated analysis
    """
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-2.5-flash')

    # Extract slide structure information
    slide_info = {
        "shapes": [],
        "text_content": []
    }

    for shape in slide.shapes:
        shape_data = {
            "type": str(shape.shape_type),
            "left": round(shape.left / 914400, 2),  # Convert to inches
            "top": round(shape.top / 914400, 2),
            "width": round(shape.width / 914400, 2),
            "height": round(shape.height / 914400, 2)
        }

        # Extract text if present
        if hasattr(shape, 'text') and shape.text.strip():
            shape_data["text"] = shape.text[:100]  # First 100 chars
            slide_info["text_content"].append(shape.text[:100])

            # Get font info
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    shape_data["font_size"] = run.font.size
                    shape_data["font_bold"] = run.font.bold
                    if run.font.color and hasattr(run.font.color, 'rgb'):
                        try:
                            shape_data["font_color"] = str(run.font.color.rgb)
                        except:
                            pass

        # Get fill color
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type == 1:  # Solid fill
                    if hasattr(shape.fill.fore_color, 'rgb'):
                        shape_data["fill_color"] = str(shape.fill.fore_color.rgb)
            except:
                pass

        slide_info["shapes"].append(shape_data)

    # Create prompt for Gemini
    prompt = f"""Analyze this PowerPoint slide structure and provide:

Slide Index: {slide_index}
Number of shapes: {len(slide_info['shapes'])}
Text content samples: {', '.join(slide_info['text_content'][:3])}

Shape details:
{json.dumps(slide_info['shapes'], indent=2)}

Please provide a JSON response with:
1. "slide_type" - What type of slide is this? (e.g., "title_slide", "section_header", "content_bullets", "two_column", "full_image", "chart_slide", etc.)
2. "usage_description" - A single concise line (max 100 chars) describing what this slide should be used for
3. "key_features" - List of 3-5 key visual/layout features
4. "recommended_content" - What type of content works best (one line)
5. "function_name" - Suggested Python function name (e.g., "add_title_slide", "add_section_header")

Return ONLY valid JSON, no other text."""

    try:
        response = model.generate_content(prompt)
        result_text = response.text.strip()

        # Clean up markdown code blocks if present
        if result_text.startswith('```'):
            lines = result_text.split('\n')
            result_text = '\n'.join(lines[1:-1])  # Remove first and last line

        result = json.loads(result_text)
        return result
    except Exception as e:
        print(f"Warning: Gemini analysis failed for slide {slide_index}: {e}")
        # Fallback to basic analysis
        return {
            "slide_type": "generic_slide",
            "usage_description": f"Slide with {len(slide_info['shapes'])} shapes",
            "key_features": ["Custom layout"],
            "recommended_content": "General content",
            "function_name": f"add_slide_{slide_index}"
        }


def generate_function_with_gemini(slide, slide_analysis: Dict[str, Any],
                                   template_id: str, api_key: str) -> str:
    """
    Use Gemini to generate Python function code for creating this slide

    Args:
        slide: Slide object
        slide_analysis: Analysis from analyze_slide_with_gemini
        template_id: Template identifier
        api_key: Google AI Studio API key

    Returns:
        Generated Python function code
    """
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-2.5-flash')

    # Extract detailed slide structure
    shapes_detail = []
    for shape in slide.shapes:
        shape_info = {
            "type": str(shape.shape_type),
            "left_inches": round(shape.left / 914400, 2),
            "top_inches": round(shape.top / 914400, 2),
            "width_inches": round(shape.width / 914400, 2),
            "height_inches": round(shape.height / 914400, 2),
            "has_text": hasattr(shape, 'text') and bool(shape.text.strip())
        }

        if shape_info["has_text"]:
            shape_info["sample_text"] = shape.text[:50]

            # Font details
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    if run.font.size:
                        shape_info["font_size_pt"] = round(run.font.size / 12700)
                    shape_info["font_bold"] = run.font.bold

                    if run.font.color and hasattr(run.font.color, 'rgb'):
                        try:
                            rgb = run.font.color.rgb
                            shape_info["font_color_rgb"] = f"RGBColor({rgb[0]}, {rgb[1]}, {rgb[2]})"
                        except:
                            pass

        # Fill color
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:
                if hasattr(shape.fill.fore_color, 'rgb'):
                    rgb = shape.fill.fore_color.rgb
                    shape_info["fill_color_rgb"] = f"RGBColor({rgb[0]}, {rgb[1]}, {rgb[2]})"
        except:
            pass

        shapes_detail.append(shape_info)

    prompt = f"""Generate a Python function to create this PowerPoint slide.

Slide Type: {slide_analysis['slide_type']}
Function Name: {slide_analysis['function_name']}
Usage: {slide_analysis['usage_description']}

Shape Details:
{json.dumps(shapes_detail, indent=2)}

Requirements:
1. Function should be a method of a class that has self.prs (Presentation object)
2. Function should clone background from self.template_slides[BEST_INDEX] where BEST_INDEX is the most appropriate template slide
3. Add textboxes on top using the exact positions from shape details
4. Use proper imports: from pptx.util import Inches, Pt; from pptx.dml.color import RGBColor
5. Parameters should accept dynamic content (title, content, bullets, etc.)
6. Include comprehensive docstring with Args description
7. Apply the exact font sizes and colors from the shape details

Return ONLY the Python function code, no markdown formatting, no explanations."""

    try:
        response = model.generate_content(prompt)
        function_code = response.text.strip()

        # Clean up markdown code blocks
        if function_code.startswith('```python'):
            lines = function_code.split('\n')
            function_code = '\n'.join(lines[1:-1])
        elif function_code.startswith('```'):
            lines = function_code.split('\n')
            function_code = '\n'.join(lines[1:-1])

        return function_code
    except Exception as e:
        print(f"Warning: Function generation failed: {e}")
        return f"""
    def {slide_analysis['function_name']}(self, title: str = "", content: str = "") -> None:
        \"\"\"
        {slide_analysis['usage_description']}

        Args:
            title: Slide title
            content: Slide content
        \"\"\"
        # Auto-generated fallback function
        slide = self._clone_slide_with_layout(0)
        pass
"""


def analyze_ppt_template_with_ai(template_file_path: str) -> Dict[str, Any]:
    """
    Enhanced template analysis using Gemini AI

    Args:
        template_file_path: Path to PPTX template

    Returns:
        Enhanced metadata with AI analysis and generated functions
    """
    if not os.path.exists(template_file_path):
        raise FileNotFoundError(f"Template file not found: {template_file_path}")

    api_key = os.getenv('GOOGLE_AI_STUDIO_KEY')
    if not api_key:
        raise ValueError("GOOGLE_AI_STUDIO_KEY not found in environment")

    prs = Presentation(template_file_path)
    template_id = generate_template_id(template_file_path)

    print(f"\n[AI Template Analyzer] Analyzing template with Gemini...")
    print(f"   Template: {os.path.basename(template_file_path)}")
    print(f"   Total slides: {len(prs.slides)}")

    metadata = {
        "template_id": template_id,
        "template_name": os.path.basename(template_file_path),
        "template_path": template_file_path,
        "slide_width_inches": prs.slide_width.inches,
        "slide_height_inches": prs.slide_height.inches,
        "total_slides": len(prs.slides),
        "powerpoint_layouts": [],
        "analyzed_slides": [],
        "generated_functions": {}
    }

    # Extract PowerPoint layouts (for reference)
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "layout_id": f"layout_{idx}",
            "layout_index": idx,
            "layout_name": layout.name,
            "placeholder_count": len(layout.placeholders)
        }
        metadata["powerpoint_layouts"].append(layout_info)

    # Analyze each slide with Gemini
    slides_list = list(prs.slides)
    generated_functions = []

    for idx, slide in enumerate(slides_list):
        print(f"[AI Template Analyzer] Analyzing slide {idx + 1}/{len(slides_list)}...")

        # Get AI analysis
        analysis = analyze_slide_with_gemini(slide, idx, api_key)

        # Generate function code
        print(f"[AI Template Analyzer] Generating function: {analysis['function_name']}...")
        function_code = generate_function_with_gemini(slide, analysis, template_id, api_key)

        slide_metadata = {
            "slide_index": idx,
            "slide_type": analysis['slide_type'],
            "usage_description": analysis['usage_description'],
            "key_features": analysis['key_features'],
            "recommended_content": analysis['recommended_content'],
            "function_name": analysis['function_name'],
            "shape_count": len(slide.shapes)
        }

        metadata["analyzed_slides"].append(slide_metadata)
        metadata["generated_functions"][analysis['function_name']] = {
            "slide_index": idx,
            "description": analysis['usage_description']
        }

        generated_functions.append({
            "function_name": analysis['function_name'],
            "code": function_code
        })

    print(f"[AI Template Analyzer] Analysis complete!")
    print(f"   Analyzed {len(metadata['analyzed_slides'])} slides")
    print(f"   Generated {len(generated_functions)} functions")

    # Store generated functions for later use
    metadata["_generated_function_codes"] = generated_functions

    return metadata


if __name__ == "__main__":
    template_path = "presentation-template.pptx"

    print("=" * 70)
    print("AI-POWERED TEMPLATE ANALYSIS")
    print("=" * 70)

    try:
        metadata = analyze_ppt_template_with_ai(template_path)

        # Save metadata
        os.makedirs("templates", exist_ok=True)
        metadata_path = f"templates/{metadata['template_id']}_metadata_ai.json"

        # Don't save function codes in JSON (too large)
        function_codes = metadata.pop("_generated_function_codes")

        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2)

        print(f"\n[PASS] Metadata saved: {metadata_path}")
        print(f"\n" + "=" * 70)
        print("ANALYZED SLIDES")
        print("=" * 70)

        for slide in metadata["analyzed_slides"]:
            print(f"\nSlide {slide['slide_index']}: {slide['slide_type']}")
            print(f"   Function: {slide['function_name']}")
            print(f"   Usage: {slide['usage_description']}")
            print(f"   Features: {', '.join(slide['key_features'])}")

    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
