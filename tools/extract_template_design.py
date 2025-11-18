"""
Extract design elements from PowerPoint template
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json


def extract_template_design(template_path: str):
    """
    Extract comprehensive design information from template

    Returns:
        Dictionary with colors, fonts, backgrounds, layouts
    """
    prs = Presentation(template_path)

    design_info = {
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "layouts": [],
        "sample_slides": []
    }

    # Extract layout information
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholders": []
        }

        for shape in layout.placeholders:
            placeholder_info = {
                "type": shape.placeholder_format.type,
                "idx": shape.placeholder_format.idx,
                "name": shape.name,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            layout_info["placeholders"].append(placeholder_info)

        design_info["layouts"].append(layout_info)

    # Extract sample slides to see actual styling
    slide_list = list(prs.slides)
    for i, slide in enumerate(slide_list[:3]):  # First 3 slides
        slide_info = {
            "slide_number": i + 1,
            "shapes": []
        }

        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }

            # Extract text if present
            if hasattr(shape, "text_frame"):
                shape_info["has_text"] = True
                shape_info["text"] = shape.text[:100] if shape.text else ""

                # Get font info from first paragraph
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        font_info = {
                            "name": run.font.name,
                            "size": run.font.size,
                            "bold": run.font.bold,
                            "italic": run.font.italic
                        }
                        if run.font.color and run.font.color.rgb:
                            font_info["color_rgb"] = str(run.font.color.rgb)
                        shape_info["font"] = font_info

            # Extract fill color
            if hasattr(shape, "fill"):
                if shape.fill.type == 1:  # Solid fill
                    if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.rgb:
                        shape_info["fill_color_rgb"] = str(shape.fill.fore_color.rgb)

            slide_info["shapes"].append(shape_info)

        design_info["sample_slides"].append(slide_info)

    return design_info


if __name__ == "__main__":
    import sys

    template_path = "presentation-template.pptx"

    print("=" * 70)
    print("TEMPLATE DESIGN EXTRACTION")
    print("=" * 70)

    design = extract_template_design(template_path)

    print(f"\nTemplate: {template_path}")
    print(f"Slides: {design['slide_count']}")
    print(f"Layouts: {design['layout_count']}")
    print(f"Dimensions: {design['slide_width']/914400:.1f}\" x {design['slide_height']/914400:.1f}\"")

    print("\n" + "=" * 70)
    print("LAYOUTS")
    print("=" * 70)
    for layout in design['layouts']:
        print(f"\nLayout {layout['index']}: {layout['name']}")
        print(f"  Placeholders: {len(layout['placeholders'])}")
        for ph in layout['placeholders']:
            print(f"    - {ph['name']} (type: {ph['type']}, idx: {ph['idx']})")

    print("\n" + "=" * 70)
    print("SAMPLE SLIDES ANALYSIS")
    print("=" * 70)
    for slide in design['sample_slides']:
        print(f"\nSlide {slide['slide_number']}:")
        print(f"  Shapes: {len(slide['shapes'])}")
        for shape in slide['shapes']:
            print(f"    - {shape['name']} ({shape['type']})")
            if 'text' in shape and shape['text']:
                print(f"      Text: {shape['text'][:50]}...")
            if 'font' in shape:
                font = shape['font']
                print(f"      Font: {font.get('name', 'N/A')}, Size: {font.get('size', 'N/A')}")
                if 'color_rgb' in font:
                    print(f"      Color: {font['color_rgb']}")
            if 'fill_color_rgb' in shape:
                print(f"      Fill: {shape['fill_color_rgb']}")

    # Save to JSON
    with open('output/template_design_analysis.json', 'w') as f:
        json.dump(design, f, indent=2, default=str)

    print("\n" + "=" * 70)
    print("Design data saved to: output/template_design_analysis.json")
    print("=" * 70)
