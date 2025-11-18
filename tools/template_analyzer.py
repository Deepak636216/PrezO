"""
Template Analysis Tool
Extracts structure from PPTX templates and generates callable functions
"""

import os
import json
import hashlib
from typing import Dict, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def generate_template_id(template_file_path: str) -> str:
    """Generate unique template ID from file path"""
    basename = os.path.basename(template_file_path)
    name_without_ext = os.path.splitext(basename)[0]
    # Create simple readable ID
    template_id = name_without_ext.lower().replace(' ', '_').replace('-', '_')
    return template_id


def analyze_ppt_template(template_file_path: str) -> Dict[str, Any]:
    """
    Extract complete structure from PPTX template

    Args:
        template_file_path: Path to uploaded PPTX

    Returns:
        Template metadata with layouts, slots, positioning
    """
    if not os.path.exists(template_file_path):
        raise FileNotFoundError(f"Template file not found: {template_file_path}")

    prs = Presentation(template_file_path)

    metadata = {
        "template_id": generate_template_id(template_file_path),
        "template_name": os.path.basename(template_file_path),
        "template_path": template_file_path,
        "slide_width_inches": prs.slide_width.inches,
        "slide_height_inches": prs.slide_height.inches,
        "total_slides": len(prs.slides),
        "layouts": [],
        "available_slide_types": []
    }

    # Extract each layout
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "layout_id": f"layout_{idx}",
            "layout_index": idx,
            "layout_name": layout.name,
            "slots": {}
        }

        # Extract placeholders
        for shape in layout.placeholders:
            slot_info = {
                "placeholder_type": str(shape.placeholder_format.type),
                "placeholder_idx": shape.placeholder_format.idx,
                "name": shape.name,
                "position": {
                    "left_inches": round(shape.left.inches, 2),
                    "top_inches": round(shape.top.inches, 2),
                    "width_inches": round(shape.width.inches, 2),
                    "height_inches": round(shape.height.inches, 2)
                }
            }

            # Try to get text info if available
            if shape.has_text_frame:
                slot_info["has_text_frame"] = True

            layout_info["slots"][shape.name] = slot_info

        # Also capture background and other shapes
        layout_info["total_shapes"] = len(layout.shapes)

        metadata["layouts"].append(layout_info)

    # Extract sample slide types from existing slides
    slide_types = []
    sample_slides = list(prs.slides)[:5]  # Analyze first 5 slides

    for idx, slide in enumerate(sample_slides):
        slide_type_info = {
            "sample_slide_index": idx,
            "description": _classify_slide_type(slide),
            "shape_count": len(slide.shapes)
        }
        slide_types.append(slide_type_info)

    metadata["available_slide_types"] = slide_types

    return metadata


def _classify_slide_type(slide) -> str:
    """Classify slide type based on content"""
    shapes = list(slide.shapes)

    # Count text shapes with content
    text_shapes = [s for s in shapes if hasattr(s, 'text') and s.text.strip()]

    if not text_shapes:
        return "Background/Design slide"

    # Check for large title text
    for shape in text_shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    font_size = para.runs[0].font.size
                    if font_size and font_size > 400000:  # Large font (>30pt)
                        if any(keyword in shape.text.lower() for keyword in ['title', 'presentation']):
                            return "Title slide"
                        else:
                            return "Section header slide"

    # Check for numbered sections
    for shape in text_shapes:
        if shape.text.strip() in ['01', '02', '03', '04', '05', '06', '07', '08', '09', '10']:
            return "Section divider slide"

    # If multiple text shapes, likely content
    if len(text_shapes) >= 3:
        return "Content slide (multi-paragraph)"

    return "Content slide"


def generate_template_functions(template_metadata: Dict[str, Any], output_dir: str = "templates") -> str:
    """
    Auto-generate Python functions from template metadata

    Args:
        template_metadata: Parsed template structure
        output_dir: Directory to save generated functions

    Returns:
        Path to generated functions.py file
    """

    template_id = template_metadata["template_id"]
    template_name = template_metadata["template_name"]

    # Create class name (capitalize and remove underscores)
    class_name = ''.join(word.capitalize() for word in template_id.split('_')) + "Template"

    # Generate code
    class_code = f'''"""
Auto-generated template functions for: {template_name}
Template ID: {template_id}
Generated by: Template Analyzer Tool
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Optional


class {class_name}:
    """
    Auto-generated template functions
    Template: {template_name}
    """

    def __init__(self, template_path: str):
        """
        Initialize template

        Args:
            template_path: Path to template PPTX file
        """
        # Load template to get layouts and styling
        self.prs = Presentation(template_path)

        # Remove all existing slides from template to start clean
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

        self.template_metadata = {json.dumps(template_metadata, indent=8)}
        self.template_path = template_path

    def add_title_slide(self, title: str, subtitle: str = None) -> None:
        """
        Add title slide

        Args:
            title: Main title text
            subtitle: Subtitle text (optional)
        """
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)

        # Add title as text box
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.2),
                Inches(8), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title: str, bullets: List[str]) -> None:
        """
        Add bullet point slide

        Args:
            title: Slide title
            bullets: List of bullet points
        """
        # Use layout 1 if available, otherwise layout 0
        layout_idx = min(1, len(self.prs.slide_layouts) - 1)
        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)

        # Add title as text box if no title placeholder
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Add bullets as text box
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(4)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = str(bullet)
            p.level = 0
            p.font.size = Pt(18)

    def add_section_header_slide(self, title: str) -> None:
        """
        Add section header slide

        Args:
            title: Section title
        """
        # Try to find section header layout (often layout 2 or 3)
        layout_idx = min(2, len(self.prs.slide_layouts) - 1)
        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

    def add_image_placeholder_slide(self, title: str, image_prompt: str,
                                     context: str = None) -> None:
        """
        Add slide with image placeholder showing prompt

        Args:
            title: Slide title
            image_prompt: AI-generated image prompt to display
            context: Optional context about why this image
        """
        # Use blank or content layout
        layout_idx = min(5, len(self.prs.slide_layouts) - 1)
        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)

        # Add title box
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Add image placeholder box
        placeholder_box = slide.shapes.add_shape(
            1,  # Rectangle (MSO_SHAPE.RECTANGLE)
            Inches(1.5), Inches(2),
            Inches(7), Inches(4)
        )

        # Style the placeholder
        placeholder_box.fill.solid()
        placeholder_box.fill.fore_color.rgb = RGBColor(232, 244, 248)  # Light blue
        placeholder_box.line.color.rgb = RGBColor(30, 58, 138)  # Dark blue
        placeholder_box.line.width = Pt(2)

        # Add prompt text inside placeholder
        text_frame = placeholder_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        text_frame.margin_top = Inches(0.3)

        # Header
        p1 = text_frame.paragraphs[0]
        p1.text = "📷 IMAGE PROMPT:"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(30, 58, 138)

        # Prompt text
        p2 = text_frame.add_paragraph()
        p2.text = image_prompt
        p2.font.size = Pt(11)
        p2.space_before = Pt(12)
        p2.font.color.rgb = RGBColor(0, 0, 0)

        # Context if provided
        if context:
            p3 = text_frame.add_paragraph()
            p3.text = f"\\n💡 Context: {{context}}"
            p3.font.size = Pt(10)
            p3.font.italic = True
            p3.space_before = Pt(12)
            p3.font.color.rgb = RGBColor(75, 85, 99)

    def add_two_column_slide(self, title: str, left_content: List[str],
                            right_content: List[str]) -> None:
        """
        Add two-column comparison slide

        Args:
            title: Slide title
            left_content: Left column bullet points
            right_content: Right column bullet points
        """
        # Use any available layout
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Add two text boxes
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(4.5), Inches(4)
        )
        right_box = slide.shapes.add_textbox(
            Inches(5), Inches(1.5),
            Inches(4.5), Inches(4)
        )

        # Fill left column
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        for item in left_content:
            p = left_frame.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(14)

        # Fill right column
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        for item in right_content:
            p = right_frame.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(14)

    def save(self, output_path: str) -> str:
        """
        Save the presentation

        Args:
            output_path: Path to save PPTX file

        Returns:
            Saved file path
        """
        self.prs.save(output_path)
        return output_path

    def get_slide_count(self) -> int:
        """Get current number of slides"""
        return len(self.prs.slides)
'''

    # Create output directory if needed
    os.makedirs(output_dir, exist_ok=True)

    # Save to file
    output_path = os.path.join(output_dir, f"{template_id}_functions.py")
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(class_code)

    # Also save metadata JSON
    metadata_path = os.path.join(output_dir, f"{template_id}_metadata.json")
    with open(metadata_path, 'w', encoding='utf-8') as f:
        json.dump(template_metadata, f, indent=2)

    return output_path


if __name__ == "__main__":
    # Test with example template
    import sys

    if len(sys.argv) > 1:
        template_path = sys.argv[1]
    else:
        print("Usage: python template_analyzer.py <path_to_template.pptx>")
        sys.exit(1)

    print(f"Analyzing template: {template_path}")
    metadata = analyze_ppt_template(template_path)

    print(f"\\nTemplate ID: {metadata['template_id']}")
    print(f"Layouts found: {len(metadata['layouts'])}")

    for layout in metadata['layouts']:
        print(f"  - {layout['layout_name']} ({len(layout['slots'])} placeholders)")

    print("\\nGenerating functions...")
    func_path = generate_template_functions(metadata)
    print(f"Functions saved to: {func_path}")
